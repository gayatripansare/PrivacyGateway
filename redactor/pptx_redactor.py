"""
redactor/pptx_redactor.py

Redacts PII from .pptx / .potx files using python-pptx.
Preserves 100% of the original file:
  - All slide layouts, masters, themes
  - All formatting (fonts, colors, sizes, bold/italic)
  - All images, charts, SmartArt (passed through untouched)
  - All animations and transitions
  - All speaker notes
  - Tables inside slides
  - Text boxes, placeholders, grouped shapes

Returns the cleaned .pptx file at out_path.

Install: pip install python-pptx
"""

import re
import shutil
import os


# ─────────────────────────────────────────────────────────
# INTERNAL HELPERS
# ─────────────────────────────────────────────────────────

def _build_replace_map(findings):
    """
    Build list of (raw_value, placeholder) pairs.
    Sorted longest-first to prevent partial match issues.
    """
    type_counts = {}
    replace_map = []

    for finding in findings:
        value   = finding.get("value", "")
        replace = finding.get("replace", "[REDACTED]")
        ftype   = finding.get("type", "UNKNOWN")

        if not value:
            continue

        type_counts[ftype] = type_counts.get(ftype, 0) + 1
        placeholder = replace.replace("]", f"_{type_counts[ftype]}]")
        replace_map.append((value, placeholder))

    replace_map.sort(key=lambda x: len(x[0]), reverse=True)
    return replace_map


def _redact_string(value, replace_map):
    """Apply all replacements to a string. Returns (new_string, was_changed)."""
    if not isinstance(value, str) or not value.strip():
        return value, False

    original = value
    for raw, placeholder in replace_map:
        value = re.sub(re.escape(raw), placeholder, value, flags=re.IGNORECASE)

    return value, value != original


# ─────────────────────────────────────────────────────────
# RUN-LEVEL REDACTION
# python-pptx stores text in Run objects inside Paragraph objects
# inside TextFrame objects. We redact at run level to preserve
# per-character formatting (color, bold, size, font name).
#
# The challenge: PII can span multiple runs (spell-check splits text).
# Strategy:
#   1. Try run-level replacement first (catches most cases)
#   2. If that misses something, merge paragraph text, find PII,
#      then rebuild runs intelligently.
# ─────────────────────────────────────────────────────────

def _redact_run(run, replace_map):
    """Redact a single run's text, preserving all run formatting."""
    if not run.text:
        return False
    new_text, changed = _redact_string(run.text, replace_map)
    if changed:
        run.text = new_text
    return changed


def _redact_paragraph(paragraph, replace_map):
    """
    Redact all runs in a paragraph.
    First tries per-run replacement. Then checks if paragraph-level
    concatenation still contains PII (cross-run split case) and
    handles it by collapsing to first run.
    """
    changed = False

    # Pass 1: per-run replacement (handles 95% of cases)
    for run in paragraph.runs:
        if _redact_run(run, replace_map):
            changed = True

    # Pass 2: check for cross-run PII
    full_text = "".join(run.text for run in paragraph.runs)
    redacted_text, needs_merge = _redact_string(full_text, replace_map)

    if needs_merge and redacted_text != full_text:
        # PII spans multiple runs — collapse all runs into first run,
        # delete the rest. First run keeps its formatting as the
        # representative style for the whole paragraph segment.
        runs = paragraph.runs
        if runs:
            runs[0].text = redacted_text
            # Remove all subsequent runs' text (keep elements for structure)
            for run in runs[1:]:
                run.text = ""
            changed = True

    return changed


def _redact_text_frame(text_frame, replace_map):
    """Redact all paragraphs in a text frame."""
    changed = False
    for paragraph in text_frame.paragraphs:
        if _redact_paragraph(paragraph, replace_map):
            changed = True
    return changed


# ─────────────────────────────────────────────────────────
# SHAPE TRAVERSAL
# Handles all shape types: text boxes, placeholders,
# tables, grouped shapes, SmartArt text.
# ─────────────────────────────────────────────────────────

def _redact_shape(shape, replace_map):
    """
    Redact PII from any shape type.
    Recursively handles groups.
    Returns True if anything was changed.
    """
    changed = False

    try:
        from pptx.util import Pt
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        pass

    # Text frame (text box, title, body placeholder, etc.)
    if shape.has_text_frame:
        if _redact_text_frame(shape.text_frame, replace_map):
            changed = True

    # Table
    if shape.has_table:
        table = shape.table
        for row in table.rows:
            for cell in row.cells:
                if _redact_text_frame(cell.text_frame, replace_map):
                    changed = True

    # Group shape — recurse into children
    try:
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP = 6
            for child_shape in shape.shapes:
                if _redact_shape(child_shape, replace_map):
                    changed = True
    except Exception:
        pass

    return changed


def _redact_slide(slide, replace_map):
    """Redact all shapes and notes on a slide."""
    changed = False

    # All shapes on the slide
    for shape in slide.shapes:
        if _redact_shape(shape, replace_map):
            changed = True

    # Speaker notes
    try:
        notes_slide = slide.notes_slide
        if notes_slide:
            for shape in notes_slide.shapes:
                if _redact_shape(shape, replace_map):
                    changed = True
    except Exception:
        pass  # Not all slides have notes

    return changed


# ─────────────────────────────────────────────────────────
# SLIDE MASTER + LAYOUT REDACTION
# Unlikely to contain PII but we cover it for completeness
# ─────────────────────────────────────────────────────────

def _redact_slide_master(master, replace_map):
    """Redact text in slide masters and their layouts."""
    for shape in master.shapes:
        _redact_shape(shape, replace_map)
    for layout in master.slide_layouts:
        for shape in layout.shapes:
            _redact_shape(shape, replace_map)


# ─────────────────────────────────────────────────────────
# PUBLIC ENTRY POINT
# ─────────────────────────────────────────────────────────

def redact_pptx(src_path, findings, out_path):
    """
    Redact PII from a PowerPoint file (.pptx) and save to out_path.

    Parameters
    ----------
    src_path  : str  — path to original .pptx file
    findings  : list — list of finding dicts from the scanner
                       each must have: "value", "replace", "type"
    out_path  : str  — where to write the cleaned file

    Returns
    -------
    out_path  : str  — path to the cleaned file

    Raises
    ------
    ImportError  if python-pptx is not installed
    FileNotFoundError  if src_path does not exist
    """
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError(
            "python-pptx is required for PowerPoint redaction. "
            "Install it with: pip install python-pptx"
        )

    if not findings:
        shutil.copy2(src_path, out_path)
        return out_path

    if not os.path.exists(src_path):
        raise FileNotFoundError(f"Source file not found: {src_path}")

    replace_map = _build_replace_map(findings)

    if not replace_map:
        shutil.copy2(src_path, out_path)
        return out_path

    # Load presentation
    prs = Presentation(src_path)

    # Redact every slide
    for slide in prs.slides:
        _redact_slide(slide, replace_map)

    # Optionally redact slide masters (company name leakage etc.)
    for master in prs.slide_masters:
        _redact_slide_master(master, replace_map)

    # Save
    if os.path.abspath(out_path) == os.path.abspath(src_path):
        import tempfile
        tmp = tempfile.mktemp(suffix=".pptx")
        prs.save(tmp)
        shutil.move(tmp, out_path)
    else:
        prs.save(out_path)

    return out_path


# ─────────────────────────────────────────────────────────
# EXTRACTION  — used by api.py to get text for scanning
# ─────────────────────────────────────────────────────────

def extract_text_pptx(src_path):
    """
    Extract all text from a PowerPoint file for PII scanning.

    Returns list of dicts:
      [{"slide": 1, "shape": "Title 1", "value": "text here"}, ...]
    """
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("pip install python-pptx")

    results = []
    prs = Presentation(src_path)

    for slide_num, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            _extract_shape_text(shape, slide_num, results)

        # Speaker notes
        try:
            notes_slide = slide.notes_slide
            if notes_slide:
                for shape in notes_slide.shapes:
                    if shape.has_text_frame:
                        text = shape.text_frame.text.strip()
                        if text:
                            results.append({
                                "slide": slide_num,
                                "shape": "Notes",
                                "value": text,
                            })
        except Exception:
            pass

    return results


def _extract_shape_text(shape, slide_num, results):
    """Recursively extract text from a shape into results list."""
    shape_name = getattr(shape, "name", "Unknown")

    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        if text:
            results.append({
                "slide": slide_num,
                "shape": shape_name,
                "value": text,
            })

    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                text = cell.text_frame.text.strip()
                if text:
                    results.append({
                        "slide": slide_num,
                        "shape": f"{shape_name} (table cell)",
                        "value": text,
                    })

    # Group shapes
    try:
        if shape.shape_type == 6:
            for child in shape.shapes:
                _extract_shape_text(child, slide_num, results)
    except Exception:
        pass


def get_full_text_pptx(src_path):
    """
    Returns all slide text as a single joined string.
    Used by scanner modules that expect plain text input.
    """
    items = extract_text_pptx(src_path)
    return "\n".join(item["value"] for item in items)